"""Docling parser plugin for document processing.

This module provides the Docling parser integration for extracting
text and structure from PDF, DOCX, PPTX, and XLSX files.
"""

import asyncio
from pathlib import Path
from typing import Any

import structlog

from src.plugins.base import (
    HealthStatus,
    ParserPlugin,
    ParsingResult,
    PluginMetadata,
    PluginType,
    SupportResult,
)

logger = structlog.get_logger(__name__)


class DoclingParser(ParserPlugin):
    """Docling parser plugin for PDF, Office, and image documents.
    
    Docling is the primary parser for structured document types
    including PDFs with text layers, Word documents, PowerPoint
    presentations, and Excel spreadsheets.
    
    Example:
        >>> parser = DoclingParser()
        >>> await parser.initialize({})
        >>> result = await parser.parse("/path/to/document.pdf")
        >>> print(result.text)
    """

    SUPPORTED_FORMATS = [
        ".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls",
        ".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".gif",
    ]

    MIME_TYPE_MAP = {
        "application/pdf": 0.95,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": 0.95,
        "application/msword": 0.90,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": 0.95,
        "application/vnd.ms-excel": 0.90,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": 0.95,
        "application/vnd.ms-powerpoint": 0.90,
        "image/jpeg": 0.80,
        "image/png": 0.80,
        "image/tiff": 0.80,
        "image/bmp": 0.70,
        "image/gif": 0.70,
    }

    def __init__(self) -> None:
        """Initialize the Docling parser."""
        self._docling_available = False
        self._document_converter = None
        self._config: dict[str, Any] = {}

    @property
    def metadata(self) -> PluginMetadata:
        """Return plugin metadata."""
        return PluginMetadata(
            id="docling",
            name="Docling Parser",
            version="1.0.0",
            type=PluginType.PARSER,
            description="Primary parser for PDF, Office docs, and images using Docling",
            author="Pipeline Team",
            supported_formats=self.SUPPORTED_FORMATS,
            requires_auth=False,
        )

    async def initialize(self, config: dict[str, Any]) -> None:
        """Initialize the parser with configuration.
        
        Args:
            config: Parser configuration options
        """
        self._config = config

        # Try to import docling
        try:
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.document import ConversionResult
            from docling.document_converter import DocumentConverter

            self._document_converter = DocumentConverter()
            self._docling_available = True
            logger.info("docling.parser.initialized", status="success")

        except ImportError as e:
            logger.warning(
                "docling.parser.import_failed",
                error=str(e),
                message="Docling not installed. Parser will use fallback methods.",
            )
            self._docling_available = False
        except Exception as e:
            logger.error(
                "docling.parser.initialization_failed",
                error=str(e),
                exc_info=True,
            )
            self._docling_available = False

    async def supports(
        self,
        file_path: str,
        mime_type: str | None = None,
    ) -> SupportResult:
        """Check if this parser supports the given file.
        
        Args:
            file_path: Path to the file
            mime_type: Optional MIME type of the file
            
        Returns:
            SupportResult indicating support status
        """
        path = Path(file_path)
        extension = path.suffix.lower()

        # Check extension
        if extension in self.SUPPORTED_FORMATS:
            confidence = 0.95
            if extension in (".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp", ".gif"):
                # Images have lower confidence for Docling
                confidence = 0.80
            return SupportResult(
                supported=True,
                confidence=confidence,
                reason=f"Supported extension: {extension}",
            )

        # Check MIME type if provided
        if mime_type and mime_type in self.MIME_TYPE_MAP:
            return SupportResult(
                supported=True,
                confidence=self.MIME_TYPE_MAP[mime_type],
                reason=f"Supported MIME type: {mime_type}",
            )

        return SupportResult(
            supported=False,
            confidence=1.0,
            reason=f"Unsupported file format: {extension}",
        )

    async def parse(
        self,
        file_path: str,
        options: dict[str, Any] | None = None,
    ) -> ParsingResult:
        """Parse a document and extract content.
        
        Args:
            file_path: Path to the file to parse
            options: Parser-specific options
            
        Returns:
            ParsingResult containing extracted content
        """
        import time

        options = options or {}
        start_time = time.time()
        
        logger.info(
            "docling.parser.parse_started",
            file_path=file_path,
            docling_available=self._docling_available,
        )

        path = Path(file_path)
        if not path.exists():
            logger.error("docling.parser.file_not_found", file_path=file_path)
            return ParsingResult(
                success=False,
                error=f"File not found: {file_path}",
            )
        
        file_size = path.stat().st_size
        file_ext = path.suffix.lower()
        logger.info(
            "docling.parser.file_info",
            file_path=file_path,
            file_size=file_size,
            file_extension=file_ext,
        )

        # Check support first
        support = await self.supports(file_path)
        if not support.supported:
            logger.warning(
                "docling.parser.unsupported_format",
                file_path=file_path,
                reason=support.reason,
            )
            return ParsingResult(
                success=False,
                error=support.reason,
            )

        try:
            # Parse based on availability
            if self._docling_available and self._document_converter:
                logger.info("docling.parser.using_docling", file_path=file_path)
                parse_result = await self._parse_with_docling(file_path, options)
            else:
                # Fallback to alternative parsing
                logger.info("docling.parser.using_fallback", file_path=file_path)
                parse_result = await self._parse_fallback(file_path, options)

            # Add processing time
            parse_result.processing_time_ms = int((time.time() - start_time) * 1000)
            
            # Log result
            if parse_result.success:
                logger.info(
                    "docling.parser.parse_success",
                    file_path=file_path,
                    parser_used=parse_result.parser_used,
                    text_length=len(parse_result.text) if parse_result.text else 0,
                    processing_time_ms=parse_result.processing_time_ms,
                )
            else:
                logger.warning(
                    "docling.parser.parse_failed",
                    file_path=file_path,
                    error=parse_result.error,
                )
            
            return parse_result

        except Exception as e:
            logger.error(
                "docling.parser.parse_exception",
                file_path=file_path,
                error=str(e),
                exc_info=True,
            )
            return ParsingResult(
                success=False,
                error=f"Parsing failed: {e!s}",
            )

    async def _parse_with_docling(
        self,
        file_path: str,
        options: dict[str, Any],
    ) -> ParsingResult:
        """Parse using Docling library.
        
        Args:
            file_path: Path to the file
            options: Parser options
            
        Returns:
            ParsingResult
        """
        logger.info("docling.parser.docling_start", file_path=file_path)

        # Convert document
        if self._document_converter is None:
            logger.error("docling.parser.converter_not_initialized")
            return ParsingResult(
                success=False,
                error="Docling converter not initialized",
            )
        
        try:
            # Run docling conversion in thread pool to avoid blocking
            # Docling's convert() is synchronous and CPU-intensive
            logger.debug("docling.parser.converting", file_path=file_path)
            result = await asyncio.to_thread(
                self._document_converter.convert,
                file_path
            )
            logger.debug("docling.parser.conversion_complete", file_path=file_path)
        except Exception as e:
            logger.error(
                "docling.parser.conversion_failed",
                file_path=file_path,
                error=str(e),
                exc_info=True,
            )
            return ParsingResult(
                success=False,
                error=f"Docling conversion failed: {e!s}",
            )

        # Extract text with error handling
        try:
            full_text = result.document.export_to_markdown()
            logger.debug(
                "docling.parser.text_extracted",
                file_path=file_path,
                text_length=len(full_text) if full_text else 0,
            )
        except Exception as e:
            logger.error(
                "docling.parser.text_extraction_failed",
                file_path=file_path,
                error=str(e),
                exc_info=True,
            )
            return ParsingResult(
                success=False,
                error=f"Text extraction failed: {e!s}",
            )

        # Extract pages (approximate from markdown sections)
        pages = self._extract_pages(result)

        # Extract metadata
        metadata = self._extract_metadata(result)

        # Extract tables
        tables = self._extract_tables(result)

        # Calculate confidence based on extraction quality
        confidence = self._calculate_confidence(result, full_text)

        # Check if we got meaningful text
        if not full_text or len(full_text.strip()) < 10:
            logger.warning(
                "docling.parser.no_text_extracted",
                file_path=file_path,
                text_length=len(full_text) if full_text else 0,
            )
            return ParsingResult(
                success=False,
                error="No text extracted from document",
            )

        return ParsingResult(
            success=True,
            text=full_text,
            pages=pages,
            metadata=metadata,
            format=Path(file_path).suffix.lower(),
            parser_used="docling",
            confidence=confidence,
            tables=tables,
        )

    async def _parse_fallback(
        self,
        file_path: str,
        options: dict[str, Any],
    ) -> ParsingResult:
        """Fallback parsing using available libraries.
        
        Args:
            file_path: Path to the file
            options: Parser options
            
        Returns:
            ParsingResult
        """
        path = Path(file_path)
        extension = path.suffix.lower()

        logger.info(
            "docling.parser.fallback_start",
            file_path=file_path,
            extension=extension,
        )

        if extension == ".pdf":
            return await self._parse_pdf_fallback(file_path)
        elif extension == ".docx":
            return await self._parse_docx_fallback(file_path)
        elif extension == ".pptx":
            return await self._parse_pptx_fallback(file_path)
        elif extension in (".txt", ".md", ".rst"):
            return await self._parse_text_fallback(file_path)
        else:
            return ParsingResult(
                success=False,
                error=f"Fallback parsing not available for {extension}. "
                      "Install docling for full support.",
            )

    async def _parse_pdf_fallback(self, file_path: str) -> ParsingResult:
        """Parse PDF using PyMuPDF as fallback.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            ParsingResult
        """
        logger.info("docling.parser.fallback_pdf_pymupdf", file_path=file_path)
        
        try:
            import fitz  # PyMuPDF  # type: ignore[import-untyped]

            def _extract_pdf():
                doc = fitz.open(file_path)
                pages: list[str] = []
                full_text = ""

                for page in doc:
                    text = page.get_text()
                    pages.append(text)
                    full_text += text + "\n"

                metadata = dict(doc.metadata) if doc.metadata else {}
                doc.close()
                
                return full_text.strip(), pages, metadata

            full_text, pages, metadata = await asyncio.to_thread(_extract_pdf)

            # Calculate simple confidence based on text extraction
            confidence = 0.8 if full_text.strip() else 0.3
            
            logger.info(
                "docling.parser.fallback_pdf_success",
                file_path=file_path,
                text_length=len(full_text),
                page_count=len(pages),
            )

            return ParsingResult(
                success=True,
                text=full_text.strip(),
                pages=pages,
                metadata=metadata,
                format=".pdf",
                parser_used="docling-fallback-pymupdf",
                confidence=confidence,
            )

        except ImportError:
            logger.error("docling.parser.pymupdf_not_available")
            return ParsingResult(
                success=False,
                error="PyMuPDF not available for PDF fallback parsing",
            )
        except Exception as e:
            logger.error(
                "docling.parser.pymupdf_failed",
                file_path=file_path,
                error=str(e),
                exc_info=True,
            )
            return ParsingResult(
                success=False,
                error=f"PyMuPDF parsing failed: {e!s}",
            )

    async def _parse_docx_fallback(self, file_path: str) -> ParsingResult:
        """Parse DOCX using python-docx as fallback.
        
        Args:
            file_path: Path to the DOCX file
            
        Returns:
            ParsingResult
        """
        logger.info("docling.parser.fallback_docx", file_path=file_path)
        
        try:
            from docx import Document  # type: ignore[import-untyped]

            def _extract_docx():
                doc = Document(file_path)
                full_text = []
                pages = []
                
                # Extract paragraphs
                for para in doc.paragraphs:
                    if para.text.strip():
                        full_text.append(para.text)
                
                # Extract tables
                for table in doc.tables:
                    table_text = []
                    for row in table.rows:
                        row_text = [cell.text for cell in row.cells]
                        table_text.append(" | ".join(row_text))
                    if table_text:
                        full_text.append("\n".join(table_text))
                
                # Try to get metadata
                metadata = {}
                if doc.core_properties:
                    props = doc.core_properties
                    if props.title:
                        metadata["title"] = props.title
                    if props.author:
                        metadata["author"] = props.author
                    if props.created:
                        metadata["created"] = str(props.created)
                
                text_content = "\n\n".join(full_text)
                # Approximate pages (rough estimate: ~500 words per page)
                if text_content:
                    words = text_content.split()
                    num_pages = max(1, len(words) // 500)
                    pages = [text_content]  # Single page for simplicity
                
                return text_content, pages, metadata

            full_text, pages, metadata = await asyncio.to_thread(_extract_docx)
            
            confidence = 0.8 if full_text.strip() else 0.3
            
            logger.info(
                "docling.parser.fallback_docx_success",
                file_path=file_path,
                text_length=len(full_text),
            )

            return ParsingResult(
                success=True,
                text=full_text,
                pages=pages,
                metadata=metadata,
                format=".docx",
                parser_used="docling-fallback-python-docx",
                confidence=confidence,
            )

        except ImportError:
            logger.error("docling.parser.python_docx_not_available")
            return ParsingResult(
                success=False,
                error="python-docx not available for DOCX fallback parsing",
            )
        except Exception as e:
            logger.error(
                "docling.parser.python_docx_failed",
                file_path=file_path,
                error=str(e),
                exc_info=True,
            )
            return ParsingResult(
                success=False,
                error=f"DOCX parsing failed: {e!s}",
            )

    async def _parse_pptx_fallback(self, file_path: str) -> ParsingResult:
        """Parse PPTX using python-pptx as fallback.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            ParsingResult
        """
        logger.info("docling.parser.fallback_pptx", file_path=file_path)
        
        try:
            from pptx import Presentation  # type: ignore[import-untyped]

            def _extract_pptx():
                prs = Presentation(file_path)
                full_text = []
                pages = []
                
                # Extract text from each slide
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = [f"--- Slide {slide_num} ---"]
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    if len(slide_text) > 1:  # More than just the header
                        slide_content = "\n".join(slide_text)
                        pages.append(slide_content)
                        full_text.append(slide_content)
                
                metadata = {
                    "slide_count": len(prs.slides),
                }
                
                return "\n\n".join(full_text), pages, metadata

            full_text, pages, metadata = await asyncio.to_thread(_extract_pptx)
            
            confidence = 0.75 if full_text.strip() else 0.3
            
            logger.info(
                "docling.parser.fallback_pptx_success",
                file_path=file_path,
                text_length=len(full_text),
                slide_count=len(pages),
            )

            return ParsingResult(
                success=True,
                text=full_text,
                pages=pages,
                metadata=metadata,
                format=".pptx",
                parser_used="docling-fallback-python-pptx",
                confidence=confidence,
            )

        except ImportError:
            logger.error("docling.parser.python_pptx_not_available")
            return ParsingResult(
                success=False,
                error="python-pptx not available for PPTX fallback parsing",
            )
        except Exception as e:
            logger.error(
                "docling.parser.python_pptx_failed",
                file_path=file_path,
                error=str(e),
                exc_info=True,
            )
            return ParsingResult(
                success=False,
                error=f"PPTX parsing failed: {e!s}",
            )

    async def _parse_text_fallback(self, file_path: str) -> ParsingResult:
        """Parse text files as fallback.
        
        Args:
            file_path: Path to the text file
            
        Returns:
            ParsingResult
        """
        logger.info("docling.parser.fallback_text", file_path=file_path)
        
        try:
            def _read_text():
                with open(file_path, encoding="utf-8", errors="ignore") as f:
                    return f.read()

            content = await asyncio.to_thread(_read_text)

            return ParsingResult(
                success=True,
                text=content,
                pages=[content],
                format=Path(file_path).suffix.lower(),
                parser_used="docling-fallback-text",
                confidence=0.95,
            )

        except Exception as e:
            logger.error(
                "docling.parser.text_fallback_failed",
                file_path=file_path,
                error=str(e),
            )
            return ParsingResult(
                success=False,
                error=f"Text parsing failed: {e}",
            )

    def _extract_pages(self, result: Any) -> list[str]:
        """Extract page texts from Docling result.
        
        Args:
            result: Docling ConversionResult
            
        Returns:
            List of page texts
        """
        pages: list[str] = []

        try:
            # Try to get pages from document
            if hasattr(result, "document") and hasattr(result.document, "pages"):
                for page in result.document.pages:
                    if hasattr(page, "text"):
                        pages.append(page.text)
                    elif hasattr(page, "export_to_text"):
                        pages.append(page.export_to_text())
        except Exception as e:
            logger.warning(f"Failed to extract pages: {e}")

        # Fallback: split markdown by headers
        if not pages:
            try:
                full_text = result.document.export_to_markdown()
                # Split by page-like separators or headers
                import re
                pages = re.split(r"\n#{1,2}\s+", full_text)
                pages = [p.strip() for p in pages if p.strip()]
            except Exception as e:
                logger.warning(f"Failed to extract pages from markdown: {e}")

        return pages

    def _extract_metadata(self, result: Any) -> dict[str, Any]:
        """Extract metadata from Docling result.
        
        Args:
            result: Docling ConversionResult
            
        Returns:
            Metadata dictionary
        """
        metadata: dict[str, Any] = {}

        try:
            if hasattr(result, "document"):
                doc = result.document

                # Try to get various metadata fields
                if hasattr(doc, "name"):
                    metadata["title"] = doc.name

                if hasattr(doc, "origin") and doc.origin:
                    origin = doc.origin
                    if hasattr(origin, "mimetype"):
                        metadata["mime_type"] = origin.mimetype
                    if hasattr(origin, "filename"):
                        metadata["filename"] = origin.filename

                # Extract document properties if available
                if hasattr(doc, "properties") and doc.properties:
                    for key, value in doc.properties.items():
                        metadata[key] = value

        except Exception as e:
            logger.warning(f"Failed to extract metadata: {e}")

        return metadata

    def _extract_tables(self, result: Any) -> list[dict[str, Any]]:
        """Extract tables from Docling result.
        
        Args:
            result: Docling ConversionResult
            
        Returns:
            List of table dictionaries
        """
        tables: list[dict[str, Any]] = []

        try:
            if hasattr(result, "document") and hasattr(result.document, "tables"):
                for table in result.document.tables:
                    table_data: dict[str, Any] = {}

                    if hasattr(table, "data"):
                        table_data["data"] = table.data

                    if hasattr(table, "export_to_dataframe"):
                        try:
                            df = table.export_to_dataframe()
                            table_data["rows"] = df.values.tolist()
                            table_data["columns"] = df.columns.tolist()
                        except Exception:
                            pass

                    tables.append(table_data)

        except Exception as e:
            logger.warning(f"Failed to extract tables: {e}")

        return tables

    def _calculate_confidence(self, result: Any, extracted_text: str) -> float:
        """Calculate confidence score for the extraction.
        
        Args:
            result: Docling ConversionResult
            extracted_text: The extracted text
            
        Returns:
            Confidence score (0.0 - 1.0)
        """
        confidence = 0.85  # Base confidence for Docling

        # Adjust based on text extraction
        if not extracted_text or len(extracted_text.strip()) < 10:
            confidence -= 0.3

        # Check for document structure
        if hasattr(result, "document"):
            doc = result.document

            # Bonus for having structure
            if hasattr(doc, "pages") and doc.pages:
                confidence += 0.05

            if hasattr(doc, "tables") and doc.tables:
                confidence += 0.05

        return min(1.0, max(0.0, confidence))

    async def health_check(self, config: dict[str, Any] | None = None) -> HealthStatus:
        """Check the health of the parser.
        
        Args:
            config: Optional configuration for health check
            
        Returns:
            HealthStatus indicating parser health
        """
        if not self._docling_available:
            return HealthStatus.DEGRADED  # pragma: no cover

        if not self._document_converter:
            return HealthStatus.UNHEALTHY

        return HealthStatus.HEALTHY
