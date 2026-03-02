#!/usr/bin/env python3
"""
Generate test files for integration testing.

This script creates sample PDF, DOCX, PPTX, and TXT files with meaningful content
for end-to-end testing of the document processing pipeline.

Usage:
    python tests/data/generate_test_files.py
"""

import os
from pathlib import Path

# Sample content about artificial intelligence (approx 600 words)
SAMPLE_CONTENT = """
Artificial Intelligence: Transforming the Modern World

Introduction

Artificial Intelligence (AI) has emerged as one of the most transformative technologies of the 21st century. From virtual assistants to autonomous vehicles, AI systems are increasingly integrated into our daily lives, reshaping industries and redefining human capabilities.

Historical Development

The concept of artificial intelligence dates back to ancient mythology, but modern AI research began in earnest during the 1950s. In 1956, the Dartmouth Conference marked the birth of AI as a formal academic discipline. Early researchers were optimistic, predicting that machines would soon match human intelligence. However, progress was slower than anticipated, leading to periods known as "AI winters" when funding and interest dwindled.

The resurgence of AI in the 2010s was driven by three key factors: the availability of massive datasets, advances in computational power through GPUs, and breakthroughs in deep learning algorithms. These developments enabled machines to achieve superhuman performance in specific tasks such as image recognition, natural language processing, and game playing.

Machine Learning Paradigms

Machine learning, a subset of AI, encompasses several approaches. Supervised learning involves training models on labeled data to make predictions. Unsupervised learning discovers hidden patterns in unlabeled data. Reinforcement learning enables agents to learn optimal behaviors through trial and error, receiving rewards or penalties for their actions.

Deep learning, which uses artificial neural networks with multiple layers, has proven particularly effective for complex pattern recognition. Convolutional neural networks excel at image processing, while recurrent neural networks and transformers have revolutionized natural language understanding.

Applications Across Industries

Healthcare has benefited enormously from AI innovations. Medical imaging systems can detect diseases with accuracy rivaling human specialists. Drug discovery has been accelerated through molecular modeling and prediction. Personalized treatment recommendations help physicians make better decisions.

In finance, AI algorithms power high-frequency trading, detect fraudulent transactions, and assess credit risk. Customer service has been transformed by chatbots and virtual assistants that handle routine inquiries 24/7. Manufacturing employs AI for quality control, predictive maintenance, and supply chain optimization.

Ethical Considerations

The rapid advancement of AI raises profound ethical questions. Bias in training data can lead to discriminatory outcomes, affecting hiring decisions, loan approvals, and criminal sentencing. Privacy concerns emerge as AI systems require vast amounts of personal data to function effectively.

The potential for job displacement creates economic anxiety, though many experts believe AI will generate new opportunities even as it automates certain tasks. Autonomous weapons and surveillance systems pose risks that society must address through careful governance.

Ensuring AI alignment with human values remains a critical challenge. Researchers work to create systems that are robust, interpretable, and beneficial. The field of AI safety seeks to prevent unintended consequences as systems become more powerful.

Future Prospects

Looking ahead, AI will likely continue its rapid evolution. General artificial intelligence, capable of performing any intellectual task that a human can, remains a distant goal but motivates ongoing research. Quantum computing may unlock new capabilities by solving problems intractable for classical computers.

Human-AI collaboration will become increasingly important. Rather than replacing humans entirely, AI tools will augment human capabilities, enabling people to focus on creative and strategic work while machines handle routine analysis and processing.

Education systems must adapt to prepare workers for an AI-driven economy. Continuous learning and skill development will be essential as job requirements evolve. Policymakers face the challenge of fostering innovation while protecting workers and ensuring broadly shared benefits.

Conclusion

Artificial intelligence represents both tremendous opportunity and significant responsibility. The technology's impact will depend on the choices we make as a society regarding development priorities, regulatory frameworks, and ethical guidelines. By approaching AI thoughtfully and inclusively, we can work toward a future where intelligent machines enhance human flourishing and address pressing global challenges.
"""

# Additional content for variety
TECH_CONTENT = """
Cloud Computing Fundamentals

Cloud computing has revolutionized how organizations deploy and manage IT infrastructure. By leveraging remote servers accessed via the internet, businesses can scale resources dynamically without capital investment in hardware.

The three primary service models are Infrastructure as a Service (IaaS), Platform as a Service (PaaS), and Software as a Service (SaaS). IaaS provides virtualized computing resources, PaaS offers development platforms, and SaaS delivers complete applications.

Major providers include Amazon Web Services, Microsoft Azure, and Google Cloud Platform. Each offers extensive service portfolios spanning compute, storage, networking, databases, machine learning, and more.

Benefits include cost reduction, global reach, rapid deployment, and automatic updates. Challenges involve security concerns, vendor lock-in risks, and managing complex multi-cloud environments.

Organizations adopting cloud strategies must consider data sovereignty, compliance requirements, and network latency. Hybrid and multi-cloud approaches allow businesses to optimize workload placement across different environments.
"""


def create_txt_file(filepath: Path, content: str) -> None:
    """Create a text file."""
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(content)
    print(f"Created: {filepath}")


def create_pdf_file(filepath: Path, content: str) -> None:
    """Create a PDF file using fpdf2."""
    try:
        from fpdf import FPDF
        
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        
        # Split content into lines and add to PDF
        lines = content.split('\n')
        for line in lines:
            if line.strip():
                # Handle Unicode by replacing problematic characters
                safe_line = line.encode('latin-1', 'replace').decode('latin-1')
                pdf.cell(200, 10, txt=safe_line, ln=True)
            else:
                pdf.ln(5)
        
        pdf.output(str(filepath))
        print(f"Created: {filepath}")
    except ImportError:
        print("fpdf2 not installed. Creating placeholder.")
        create_txt_file(filepath.with_suffix('.txt'), content)


def create_docx_file(filepath: Path, content: str) -> None:
    """Create a DOCX file using python-docx."""
    try:
        from docx import Document
        from docx.shared import Pt
        
        doc = Document()
        
        # Split content and add paragraphs
        lines = content.split('\n')
        for line in lines:
            if line.strip():
                p = doc.add_paragraph(line.strip())
            else:
                doc.add_paragraph()
        
        doc.save(str(filepath))
        print(f"Created: {filepath}")
    except ImportError:
        print("python-docx not installed. Creating placeholder.")
        create_txt_file(filepath.with_suffix('.txt'), content)


def create_pptx_file(filepath: Path, content: str) -> None:
    """Create a PPTX file using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        prs = Presentation()
        
        # Split content into sections for different slides
        sections = content.split('\n\n')
        
        for section in sections[:10]:  # Limit to 10 slides
            if not section.strip():
                continue
                
            slide_layout = prs.slide_layouts[1]  # Title and content
            slide = prs.slides.add_slide(slide_layout)
            
            lines = section.strip().split('\n')
            if lines:
                # First line as title
                slide.shapes.title.text = lines[0][:50]  # Limit title length
                
                # Rest as content
                if len(lines) > 1:
                    body_shape = slide.placeholders[1]
                    tf = body_shape.text_frame
                    tf.text = '\n'.join(lines[1:])
        
        prs.save(str(filepath))
        print(f"Created: {filepath}")
    except ImportError:
        print("python-pptx not installed. Creating placeholder.")
        create_txt_file(filepath.with_suffix('.txt'), content)


def main():
    """Generate all test files."""
    # Get the directory where this script is located
    script_dir = Path(__file__).parent
    
    print(f"Generating test files in: {script_dir}")
    print("-" * 50)
    
    # Create sample.pdf
    create_pdf_file(script_dir / "sample.pdf", SAMPLE_CONTENT)
    
    # Create sample.docx
    create_docx_file(script_dir / "sample.docx", SAMPLE_CONTENT)
    
    # Create sample.pptx
    create_pptx_file(script_dir / "sample.pptx", SAMPLE_CONTENT)
    
    # Create sample.txt
    create_txt_file(script_dir / "sample.txt", SAMPLE_CONTENT)
    
    # Create tech_report.pdf (different content)
    create_pdf_file(script_dir / "tech_report.pdf", TECH_CONTENT)
    
    # Create additional_test.docx (different content)
    create_docx_file(script_dir / "additional_test.docx", TECH_CONTENT)
    
    print("-" * 50)
    print("Test file generation complete!")
    
    # List created files
    print("\nGenerated files:")
    for f in script_dir.iterdir():
        if f.is_file() and f.suffix in ['.pdf', '.docx', '.pptx', '.txt']:
            size = f.stat().st_size
            print(f"  {f.name}: {size:,} bytes")


if __name__ == "__main__":
    main()
