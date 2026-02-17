#!/usr/bin/env python3
"""Generate test documents for E2E testing.

This script generates various test documents for the E2E test suite:
- PDF files (text-based, scanned/image-based, mixed, large)
- Office documents (Word, Excel, PowerPoint)
- Images (receipts, documents)
- Archives (ZIP with multiple documents)

Usage:
    python generate_test_docs.py

Requirements:
    pip install reportlab openpyxl python-pptx Pillow pypdf
"""

import json
import zipfile
from datetime import datetime
from pathlib import Path


def ensure_directories():
    """Create necessary directories."""
    base_path = Path(__file__).parent
    for subdir in ["pdf", "office", "images", "archives"]:
        (base_path / subdir).mkdir(parents=True, exist_ok=True)


def generate_text_pdf():
    """Generate a text-based PDF with extractable text."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
    except ImportError:
        print("reportlab not installed. Skipping PDF generation.")
        return

    output_path = Path(__file__).parent / "pdf" / "sample-text.pdf"

    c = canvas.Canvas(str(output_path), pagesize=letter)
    width, height = letter

    # Title
    c.setFont("Helvetica-Bold", 24)
    c.drawString(72, height - 72, "Sample Text Document")

    # Content
    c.setFont("Helvetica", 12)
    text = """
    This is a sample text-based PDF document for E2E testing.
    
    Lorem ipsum dolor sit amet, consectetur adipiscing elit. Sed do
    eiusmod tempor incididunt ut labore et dolore magna aliqua. Ut enim
    ad minim veniam, quis nostrud exercitation ullamco laboris nisi ut
    aliquip ex ea commodo consequat.
    
    Key Features:
    - Extractable text content
    - Standard PDF format
    - No images or complex formatting
    - Suitable for Docling parser
    
    Test Metadata:
    - Test Case: TC001
    - Expected Parser: Docling
    - Processing Time: ~5 seconds
    """

    y = height - 120
    for line in text.strip().split("\n"):
        c.drawString(72, y, line.strip())
        y -= 14

    c.save()
    print(f"Generated: {output_path}")


def generate_scanned_pdf():
    """Generate an image-based PDF (scanned document simulation)."""
    try:
        from PIL import Image, ImageDraw, ImageFont
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
    except ImportError:
        print("reportlab or Pillow not installed. Skipping scanned PDF generation.")
        return

    output_path = Path(__file__).parent / "pdf" / "sample-scanned.pdf"

    # Create an image to embed in PDF
    img = Image.new("RGB", (612, 792), color="white")
    draw = ImageDraw.Draw(img)

    # Add text as image
    try:
        font_large = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 32)
        font_normal = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
    except:
        font_large = ImageFont.load_default()
        font_normal = ImageFont.load_default()

    # Add some noise to simulate scanned document
    import random
    for _ in range(1000):
        x, y = random.randint(0, 612), random.randint(0, 792)
        draw.point((x, y), fill=(200, 200, 200))

    draw.text((50, 50), "Scanned Document Sample", fill="black", font=font_large)

    lines = [
        "This document simulates a scanned PDF.",
        "The text is embedded as an image and requires OCR.",
        "",
        "Expected processing:",
        "- Azure OCR should be used",
        "- Processing time: ~15 seconds",
        "- Test Case: TC002"
    ]

    y = 120
    for line in lines:
        draw.text((50, y), line, fill="black", font=font_normal)
        y += 25

    # Save image temporarily
    temp_img = Path(__file__).parent / "temp_scanned.png"
    img.save(temp_img)

    # Create PDF with image
    c = canvas.Canvas(str(output_path), pagesize=letter)
    c.drawImage(str(temp_img), 0, 0, width=612, height=792)
    c.save()

    # Cleanup
    temp_img.unlink()
    print(f"Generated: {output_path}")


def generate_mixed_pdf():
    """Generate a PDF with both text and images."""
    try:
        from PIL import Image, ImageDraw
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.utils import ImageReader
        from reportlab.pdfgen import canvas
    except ImportError:
        print("reportlab or Pillow not installed. Skipping mixed PDF generation.")
        return

    output_path = Path(__file__).parent / "pdf" / "sample-mixed.pdf"

    c = canvas.Canvas(str(output_path), pagesize=letter)
    width, height = letter

    # Add text
    c.setFont("Helvetica-Bold", 20)
    c.drawString(72, height - 72, "Mixed Content PDF")

    c.setFont("Helvetica", 12)
    c.drawString(72, height - 100, "This PDF contains both text and images.")

    # Add a simple image
    img = Image.new("RGB", (200, 100), color="lightblue")
    draw = ImageDraw.Draw(img)
    draw.rectangle([10, 10, 190, 90], outline="black", width=2)
    draw.text((60, 40), "Sample Image", fill="black")

    # Embed image in PDF
    img_reader = ImageReader(img)
    c.drawImage(img_reader, 72, height - 250, width=200, height=100)

    # More text
    c.drawString(72, height - 280, "More text content after the image.")
    c.drawString(72, height - 300, "Test Case: TC003")

    c.save()
    print(f"Generated: {output_path}")


def generate_large_pdf():
    """Generate a large PDF document (50+ pages) for performance testing."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
    except ImportError:
        print("reportlab not installed. Skipping large PDF generation.")
        return

    output_path = Path(__file__).parent / "pdf" / "large-document.pdf"

    c = canvas.Canvas(str(output_path), pagesize=letter)
    width, height = letter

    num_pages = 55

    for page_num in range(1, num_pages + 1):
        # Header
        c.setFont("Helvetica-Bold", 16)
        c.drawString(72, height - 72, f"Large Document - Page {page_num} of {num_pages}")

        # Content
        c.setFont("Helvetica", 10)

        # Generate some content for each page
        text = f"""
        This is page {page_num} of a large document for performance testing.
        
        Lorem ipsum dolor sit amet, consectetur adipiscing elit. Sed do eiusmod 
        tempor incididunt ut labore et dolore magna aliqua. Ut enim ad minim veniam, 
        quis nostrud exercitation ullamco laboris nisi ut aliquip ex ea commodo 
        consequat. Duis aute irure dolor in reprehenderit in voluptate velit esse 
        cillum dolore eu fugiat nulla pariatur.
        
        Section {page_num}.1: Introduction
        
        Sed ut perspiciatis unde omnis iste natus error sit voluptatem accusantium 
        doloremque laudantium, totam rem aperiam, eaque ipsa quae ab illo inventore 
        veritatis et quasi architecto beatae vitae dicta sunt explicabo.
        
        Section {page_num}.2: Details
        
        Nemo enim ipsam voluptatem quia voluptas sit aspernatur aut odit aut fugit, 
        sed quia consequuntur magni dolores eos qui ratione voluptatem sequi nesciunt.
        
        Test Case: TC004
        Expected Duration: ~60 seconds
        """

        y = height - 100
        for line in text.strip().split("\n"):
            c.drawString(72, y, line.strip())
            y -= 12

        # Footer
        c.setFont("Helvetica", 8)
        c.drawString(72, 50, f"Page {page_num}")

        if page_num < num_pages:
            c.showPage()

    c.save()
    print(f"Generated: {output_path} ({num_pages} pages)")


def generate_word_doc():
    """Generate a Word document."""
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Inches, Pt
    except ImportError:
        print("python-docx not installed. Skipping Word document generation.")
        return

    output_path = Path(__file__).parent / "office" / "sample.docx"

    doc = Document()

    # Title
    title = doc.add_heading("Sample Word Document", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Introduction
    doc.add_heading("Introduction", level=1)
    doc.add_paragraph(
        "This is a sample Word document for E2E testing. "
        "It contains formatted text, headings, and other elements."
    )

    # Features section
    doc.add_heading("Features", level=1)
    features = doc.add_paragraph()
    features.add_run("Key features include:\n").bold = True
    features.add_run("• Formatted text with bold and italic\n")
    features.add_run("• Multiple heading levels\n")
    features.add_run("• Bullet points and lists\n")
    features.add_run("• Tables and other elements\n")

    # Add a table
    doc.add_heading("Sample Table", level=2)
    table = doc.add_table(rows=3, cols=3)
    table.style = "Light Grid Accent 1"

    # Header row
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = "Column 1"
    hdr_cells[1].text = "Column 2"
    hdr_cells[2].text = "Column 3"

    # Data rows
    for i in range(1, 3):
        row_cells = table.rows[i].cells
        row_cells[0].text = f"Row {i}, Cell 1"
        row_cells[1].text = f"Row {i}, Cell 2"
        row_cells[2].text = f"Row {i}, Cell 3"

    # Test metadata
    doc.add_heading("Test Metadata", level=1)
    doc.add_paragraph("Test Case: TC005")
    doc.add_paragraph("Expected Parser: Docling")
    doc.add_paragraph("Expected Duration: ~3 seconds")

    doc.save(output_path)
    print(f"Generated: {output_path}")


def generate_excel_doc():
    """Generate an Excel spreadsheet."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font, PatternFill
    except ImportError:
        print("openpyxl not installed. Skipping Excel document generation.")
        return

    output_path = Path(__file__).parent / "office" / "sample.xlsx"

    wb = Workbook()
    ws = wb.active
    ws.title = "Sample Data"

    # Header row
    headers = ["ID", "Name", "Category", "Value", "Date"]
    for col, header in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        cell.alignment = Alignment(horizontal="center")

    # Sample data
    sample_data = [
        [1, "Item A", "Category 1", 100.50, "2024-01-15"],
        [2, "Item B", "Category 2", 200.75, "2024-01-16"],
        [3, "Item C", "Category 1", 150.00, "2024-01-17"],
        [4, "Item D", "Category 3", 300.25, "2024-01-18"],
        [5, "Item E", "Category 2", 175.50, "2024-01-19"],
    ]

    for row_idx, row_data in enumerate(sample_data, 2):
        for col_idx, value in enumerate(row_data, 1):
            ws.cell(row=row_idx, column=col_idx, value=value)

    # Add summary row
    summary_row = len(sample_data) + 2
    ws.cell(row=summary_row, column=1, value="Total").font = Font(bold=True)
    ws.cell(row=summary_row, column=4, value="=SUM(D2:D6)").font = Font(bold=True)

    # Adjust column widths
    for col in range(1, 6):
        ws.column_dimensions[chr(64 + col)].width = 15

    # Add metadata sheet
    meta_ws = wb.create_sheet(title="Metadata")
    meta_ws.cell(row=1, column=1, value="Test Case")
    meta_ws.cell(row=1, column=2, value="TC006")
    meta_ws.cell(row=2, column=1, value="Expected Parser")
    meta_ws.cell(row=2, column=2, value="Docling")
    meta_ws.cell(row=3, column=1, value="Expected Duration")
    meta_ws.cell(row=3, column=2, value="~3 seconds")

    wb.save(output_path)
    print(f"Generated: {output_path}")


def generate_powerpoint_doc():
    """Generate a PowerPoint presentation."""
    try:
        from pptx import Presentation
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError:
        print("python-pptx not installed. Skipping PowerPoint generation.")
        return

    output_path = Path(__file__).parent / "office" / "sample.pptx"

    prs = Presentation()

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    slide1 = prs.slides.add_slide(title_slide_layout)
    slide1.shapes.title.text = "Sample PowerPoint Presentation"
    slide1.placeholders[1].text = "E2E Test Document\n\nTest Case: TC007"

    # Slide 2: Content
    bullet_slide_layout = prs.slide_layouts[1]  # Title and Content
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    slide2.shapes.title.text = "Key Features"

    content = slide2.placeholders[1]
    tf = content.text_frame
    tf.text = "Document Processing Pipeline"

    p = tf.add_paragraph()
    p.text = "Multiple file format support"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Intelligent parser selection"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "OCR for scanned documents"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Destination-agnostic output"
    p.level = 1

    # Slide 3: Summary
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    slide3.shapes.title.text = "Test Information"

    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = "Processing Details:"

    p = tf3.add_paragraph()
    p.text = "Expected Parser: Docling"
    p.level = 1

    p = tf3.add_paragraph()
    p.text = "Expected Duration: ~5 seconds"
    p.level = 1

    p = tf3.add_paragraph()
    p.text = "Slide Count: 3"
    p.level = 1

    prs.save(output_path)
    print(f"Generated: {output_path}")


def generate_receipt_image():
    """Generate a receipt image for OCR testing."""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("Pillow not installed. Skipping receipt image generation.")
        return

    output_path = Path(__file__).parent / "images" / "receipt.jpg"

    # Create receipt image
    width, height = 400, 600
    img = Image.new("RGB", (width, height), color="white")
    draw = ImageDraw.Draw(img)

    try:
        font_title = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 20)
        font_normal = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 14)
        font_small = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 12)
    except:
        font_title = ImageFont.load_default()
        font_normal = ImageFont.load_default()
        font_small = ImageFont.load_default()

    # Add some noise
    import random
    for _ in range(500):
        x, y = random.randint(0, width), random.randint(0, height)
        draw.point((x, y), fill=(220, 220, 220))

    # Store info
    draw.text((width//2 - 80, 30), "SAMPLE STORE", fill="black", font=font_title)
    draw.text((width//2 - 60, 60), "123 Main Street", fill="black", font=font_small)
    draw.text((width//2 - 50, 80), "City, ST 12345", fill="black", font=font_small)

    # Date and time
    draw.text((50, 120), f"Date: {datetime.now().strftime('%Y-%m-%d')}", fill="black", font=font_normal)
    draw.text((50, 145), f"Time: {datetime.now().strftime('%H:%M:%S')}", fill="black", font=font_normal)

    # Items
    y = 190
    draw.line([(30, y-10), (width-30, y-10)], fill="black", width=1)

    items = [
        ("Item 1 Description", 19.99),
        ("Item 2 Description", 29.99),
        ("Item 3 Description", 9.99),
    ]

    for item, price in items:
        draw.text((50, y), item, fill="black", font=font_normal)
        draw.text((width-100, y), f"${price:.2f}", fill="black", font=font_normal)
        y += 30

    draw.line([(30, y+5), (width-30, y+5)], fill="black", width=1)

    # Totals
    y += 25
    draw.text((50, y), "Subtotal:", fill="black", font=font_normal)
    draw.text((width-100, y), "$59.97", fill="black", font=font_normal)

    y += 25
    draw.text((50, y), "Tax (8%):", fill="black", font=font_normal)
    draw.text((width-100, y), "$4.80", fill="black", font=font_normal)

    y += 30
    draw.text((50, y), "TOTAL:", fill="black", font=font_title)
    draw.text((width-120, y), "$64.77", fill="black", font=font_title)

    # Footer
    draw.text((width//2 - 70, height-80), "Thank you for your business!", fill="black", font=font_small)
    draw.text((width//2 - 50, height-50), "Test Case: TC008", fill="black", font=font_small)

    img.save(output_path, quality=90)
    print(f"Generated: {output_path}")


def generate_document_image():
    """Generate a clean document image."""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("Pillow not installed. Skipping document image generation.")
        return

    output_path = Path(__file__).parent / "images" / "document.png"

    # Create clean document image
    width, height = 800, 1000
    img = Image.new("RGB", (width, height), color="white")
    draw = ImageDraw.Draw(img)

    try:
        font_title = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28)
        font_header = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 18)
        font_normal = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 14)
    except:
        font_title = ImageFont.load_default()
        font_header = ImageFont.load_default()
        font_normal = ImageFont.load_default()

    # Title
    draw.text((50, 50), "Sample Document Image", fill="black", font=font_title)

    # Content
    y = 120
    draw.text((50, y), "Introduction", fill="black", font=font_header)
    y += 40

    paragraphs = [
        "This is a clean document image for OCR testing. The text is",
        "clearly visible and should be easily extractable by OCR engines.",
        "",
        "Key Information:",
        "• Document ID: DOC-2024-001",
        "• Date: January 15, 2024",
        "• Author: Test System",
        "• Classification: Public",
        "",
        "This document contains standard text without any special formatting",
        "or complex layouts. It is designed to test basic OCR capabilities.",
        "",
        "Expected Results:",
        "- Text extraction accuracy: >95%",
        "- Processing time: <5 seconds",
        "- Quality score: >0.8",
        "",
        "Test Case: TC009",
        "Expected Parser: Azure OCR"
    ]

    for line in paragraphs:
        draw.text((50, y), line, fill="black", font=font_normal)
        y += 25

    # Add border
    draw.rectangle([(10, 10), (width-10, height-10)], outline="gray", width=2)

    img.save(output_path)
    print(f"Generated: {output_path}")


def generate_archive():
    """Generate a ZIP archive with multiple documents."""
    output_path = Path(__file__).parent / "archives" / "documents.zip"

    # Ensure sample files exist before creating archive
    base_path = Path(__file__).parent
    files_to_archive = [
        (base_path / "pdf" / "sample-text.pdf", "text-document.pdf"),
        (base_path / "office" / "sample.docx", "word-document.docx"),
    ]

    with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for file_path, arcname in files_to_archive:
            if file_path.exists():
                zf.write(file_path, arcname)
                print(f"  Added to archive: {arcname}")
            else:
                # Create placeholder if file doesn't exist
                zf.writestr(arcname + ".placeholder", f"Placeholder for {arcname}")
                print(f"  Added placeholder: {arcname}")

        # Add manifest
        manifest = {
            "archive_id": "TC010",
            "description": "Test archive for bulk processing",
            "created_at": datetime.now().isoformat(),
            "files": [arcname for _, arcname in files_to_archive]
        }
        zf.writestr("manifest.json", json.dumps(manifest, indent=2))

    print(f"Generated: {output_path}")


def main():
    """Generate all test documents."""
    print("=" * 60)
    print("Generating E2E Test Documents")
    print("=" * 60)

    ensure_directories()

    print("\n--- PDF Documents ---")
    generate_text_pdf()
    generate_scanned_pdf()
    generate_mixed_pdf()
    generate_large_pdf()

    print("\n--- Office Documents ---")
    generate_word_doc()
    generate_excel_doc()
    generate_powerpoint_doc()

    print("\n--- Images ---")
    generate_receipt_image()
    generate_document_image()

    print("\n--- Archives ---")
    generate_archive()

    print("\n" + "=" * 60)
    print("Test document generation complete!")
    print("=" * 60)


if __name__ == "__main__":
    main()
